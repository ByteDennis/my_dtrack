#!/usr/bin/env python3
"""Generate PCDS Oracle Query Performance Optimization PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ---------------------------------------------------------------------------
# Theme colors
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_DARK = RGBColor(0x1B, 0x2A, 0x4A)       # dark navy
BG_LIGHT = RGBColor(0xF5, 0xF6, 0xFA)      # off-white
ACCENT = RGBColor(0x27, 0x7D, 0xA1)        # teal
ACCENT2 = RGBColor(0xE8, 0x6F, 0x51)       # coral/red for highlights
TEXT_DARK = RGBColor(0x23, 0x23, 0x23)
TEXT_MED = RGBColor(0x55, 0x55, 0x55)
TEXT_LIGHT = RGBColor(0x99, 0x99, 0x99)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GRAY = RGBColor(0xBD, 0xBD, 0xBD)
CODE_BG = RGBColor(0xF0, 0xF0, 0xF0)

FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=12,
                  bold=False, color=TEXT_DARK, font_name=FONT_BODY,
                  alignment=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_para(text_frame, text, font_size=12, bold=False, color=TEXT_DARK,
              font_name=FONT_BODY, alignment=PP_ALIGN.LEFT, space_before=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p


def _add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Data
# ---------------------------------------------------------------------------

STRATEGIES = [
    {"num": 1, "name": "Date Range WHERE Filter",
     "desc": "TRUNC(date) vs date range",
     "before": "80 s", "after": "0.05 s", "speedup": "1600x", "sas": "81 -> 0.05 s",
     "speedup_val": 1600,
     "goal": "Speed up date filters by using range predicates instead of TRUNC(), which disables index/partition usage on the date column.",
     "outcome": "80 s -> 0.05 s (1600x). Range condition enabled index on EFF_DT and partition pruning. Biggest single improvement.",
     "takeaway": "Never wrap indexed date columns in functions inside WHERE clauses.",
     "snippet": "[SNIPPET — paste SQL here]"},
    {"num": 2, "name": "Direct GROUP BY",
     "desc": "TRUNC(date) GROUP BY vs direct column",
     "before": "7.9 s", "after": "2.3 s", "speedup": "3.5x", "sas": "8 -> 2.3 s",
     "speedup_val": 3.5,
     "goal": "Group by the date column directly instead of TRUNC(date_col) to allow index usage and avoid per-row function overhead.",
     "outcome": "7.9 s -> 2.3 s (3.5x). Direct grouping utilized index on DW_BUS_DT. Modest but consistent improvement.",
     "takeaway": "If data is already at the needed granularity, group by the column directly.",
     "snippet": "[SNIPPET — paste SQL here]"},
    {"num": 3, "name": "Skip Unnecessary CTE",
     "desc": "CTE overhead vs direct query",
     "before": "11 s", "after": "1.3 s", "speedup": "8x", "sas": "12 -> 1.3 s",
     "speedup_val": 8,
     "goal": "Eliminate overhead of trivial CTEs. For simple one-off aggregations, a direct query avoids extra parsing and execution steps.",
     "outcome": "11 s -> 1.3 s (8x). The CTE wrapper added overhead without benefit for a single-use aggregation.",
     "takeaway": "Use direct queries for simple aggregations; reserve CTEs for reuse.",
     "snippet": "[SNIPPET — paste SQL here]"},
    {"num": 4, "name": "Parallel Query (PARALLEL hint)",
     "desc": "Serial vs PARALLEL(4)",
     "before": "2.1 s", "after": "3.9 s", "speedup": "none", "sas": "no gain",
     "speedup_val": 1,
     "goal": "Accelerate large aggregations via Oracle parallel execution with PARALLEL(table, 4) hint.",
     "outcome": "No gain (2.1 s serial vs 3.9 s parallel). SESSIONS_PER_USER=3 limited available threads.",
     "takeaway": "Parallel hints require sufficient session slots. Check DBA limits first.",
     "snippet": "[SNIPPET — paste SQL here]"},
    {"num": 5, "name": "Client Fetch Tuning (ARRAYSIZE)",
     "desc": "arraysize 100 vs 5000",
     "before": "20.8 s", "after": "1.3 s", "speedup": "16x", "sas": "1.5 -> 1.3 s",
     "speedup_val": 16,
     "goal": "Reduce network round-trips for large result sets by increasing Python cursor arraysize from 100 (default) to 5000.",
     "outcome": "20.8 s -> 1.3 s (16x) for 200k rows. Fewer round-trips between client and DB.",
     "takeaway": "Always set arraysize >= 5000 for bulk data retrieval in Python.",
     "snippet": "[SNIPPET — paste Python here]"},
    {"num": 6, "name": "Bind Variables",
     "desc": "Literals vs bind parameters",
     "before": "(est.)", "after": "(est.)", "speedup": "~10x", "sas": "auto binds",
     "speedup_val": 10,
     "goal": "Reduce SQL parsing overhead for repetitive lookups by using bind parameters instead of literal values in looped queries.",
     "outcome": "~10x estimated improvement for 500 single-row lookups. Bind parameters allow Oracle to reuse execution plans.",
     "takeaway": "For batch/looped queries, always use bind parameters.",
     "snippet": "[SNIPPET — paste SQL here]"},
    {"num": 7, "name": "Partition Pruning",
     "desc": "TO_CHAR vs date range on partition key",
     "before": "54 s", "after": "0.3 s", "speedup": "180x", "sas": "60 -> 0.3 s",
     "speedup_val": 180,
     "goal": "Enable Oracle partition pruning by filtering with native date range instead of TO_CHAR() on the partition key.",
     "outcome": "54 s -> 0.3 s (180x). TO_CHAR forced full-partition scan; date range restricted scan to Jan 2024 partition only.",
     "takeaway": "Never apply functions to partition key columns in WHERE clauses.",
     "snippet": "[SNIPPET — paste SQL here]"},
    {"num": 8, "name": "RESULT_CACHE Hint",
     "desc": "With vs without RESULT_CACHE",
     "before": "10.5 s", "after": "0.47 s", "speedup": "22x", "sas": "N/A",
     "speedup_val": 22,
     "goal": "Leverage Oracle's server-side result cache for repeated identical queries using the /*+ RESULT_CACHE */ hint.",
     "outcome": "10.5 s -> 0.47 s avg across 5 runs (22x). Subsequent runs return cached results sub-second.",
     "takeaway": "Use RESULT_CACHE for read-heavy, repeatedly executed queries.",
     "snippet": "[SNIPPET — paste SQL here]"},
    {"num": 9, "name": "Batch Fetching (fetchmany)",
     "desc": "fetchone vs fetchmany(5000)",
     "before": "(est.)", "after": "(est.)", "speedup": "~15x", "sas": "bulk fetch",
     "speedup_val": 15,
     "goal": "Minimize application-level overhead by fetching rows in batches of 5000 instead of one-by-one with fetchone().",
     "outcome": "~15x estimated improvement for 100k rows. Reduces Python-to-DB call count dramatically.",
     "takeaway": "Use fetchmany() or set large arraysize; never loop with fetchone().",
     "snippet": "[SNIPPET — paste Python here]"},
]

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, BG_DARK)

    # Title
    _add_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(1),
                  "PCDS Oracle Query Performance\nOptimization",
                  font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    _add_text_box(slide, Inches(1), Inches(3.3), Inches(8), Inches(0.6),
                  "Summary of Findings — SQL Strategy Benchmarks",
                  font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # Footer
    _add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.4),
                  "Tested via python-oracledb  |  SAS/ACCESS reference timings  |  5 iterations averaged",
                  font_size=10, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    # Accent line
    _add_shape_rect(slide, Inches(3.5), Inches(3.15), Inches(3), Pt(3), ACCENT)


def build_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_LIGHT)

    # Title
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                  "Executive Summary", font_size=24, bold=True, color=BG_DARK)

    # Subheading
    _add_text_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(0.4),
                  "Key Result: 10x–20x faster queries, some over 100x",
                  font_size=14, bold=True, color=ACCENT2)

    # Table
    rows = len(STRATEGIES) + 1
    cols = 6
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(0.3), Inches(1.3),
                                        Inches(9.4), Inches(4.5))
    tbl = tbl_shape.table

    # Column widths
    col_widths = [Inches(0.4), Inches(2.8), Inches(1.2), Inches(1.2), Inches(1.3), Inches(2.5)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    # Header row
    headers = ["#", "Strategy", "Before", "After", "Speed-Up", "SAS Reference"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_DARK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    highlight_rows = {1, 7, 8}  # highest gains
    nogain_rows = {4}
    for i, s in enumerate(STRATEGIES):
        row_idx = i + 1
        vals = [str(s["num"]), s["name"], s["before"], s["after"], s["speedup"], s["sas"]]
        for j, v in enumerate(vals):
            cell = tbl.cell(row_idx, j)
            cell.text = v
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = FONT_BODY
                p.alignment = PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT
                if s["num"] in highlight_rows:
                    p.font.color.rgb = ACCENT2 if j == 4 else TEXT_DARK
                    p.font.bold = (j == 4)
                elif s["num"] in nogain_rows:
                    p.font.color.rgb = TEXT_LIGHT
                else:
                    p.font.color.rgb = TEXT_DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEB, 0xEE, 0xF5)

    # Row heights
    tbl.rows[0].height = Inches(0.4)
    for i in range(1, rows):
        tbl.rows[i].height = Inches(0.42)


def build_chart_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_LIGHT)

    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                  "Speed-Up Factor by Strategy", font_size=24, bold=True, color=BG_DARK)
    _add_text_box(slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.3),
                  "Logarithmic scale — higher is better. Red dashed line = 20x reference.",
                  font_size=10, color=TEXT_MED)

    # Sort by speedup descending
    sorted_strats = sorted(STRATEGIES, key=lambda s: s["speedup_val"], reverse=True)

    chart_data = CategoryChartData()
    chart_data.categories = [f"#{s['num']}. {s['name']}" for s in sorted_strats]
    chart_data.add_series("Speed-Up (x)", [s["speedup_val"] for s in sorted_strats])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.2),
        Inches(9), Inches(5.0), chart_data)

    chart = chart_frame.chart
    chart.has_legend = False

    # Style bars
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT

    # Color individual points
    for i, s in enumerate(sorted_strats):
        pt = series.points[i]
        pt.format.fill.solid()
        if s["speedup_val"] >= 100:
            pt.format.fill.fore_color.rgb = ACCENT2
        elif s["speedup_val"] <= 1:
            pt.format.fill.fore_color.rgb = GRAY
        else:
            pt.format.fill.fore_color.rgb = ACCENT

    # Value axis (horizontal for bar chart)
    val_axis = chart.value_axis
    val_axis.has_title = False
    val_axis.minimum_scale = 0
    val_axis.maximum_scale = 1800
    val_axis.major_gridlines.format.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    # Category axis
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.name = FONT_BODY

    # Data labels
    series.has_data_labels = True
    series.data_labels.font.size = Pt(9)
    series.data_labels.font.bold = True
    series.data_labels.font.color.rgb = TEXT_DARK
    series.data_labels.number_format = '#,##0"x"'
    series.data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END


def build_strategy_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    # Top accent bar
    _add_shape_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), ACCENT)

    # Number badge
    badge = _add_shape_rect(slide, Inches(0.4), Inches(0.3), Inches(0.5), Inches(0.5), BG_DARK)
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(s["num"])
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Title
    _add_text_box(slide, Inches(1.1), Inches(0.3), Inches(8), Inches(0.5),
                  s["name"], font_size=22, bold=True, color=BG_DARK)

    # Subtitle / description
    _add_text_box(slide, Inches(1.1), Inches(0.75), Inches(5), Inches(0.3),
                  s["desc"], font_size=11, color=TEXT_MED)

    # LEFT COLUMN: Goal + Snippet
    _add_text_box(slide, Inches(0.4), Inches(1.3), Inches(5.5), Inches(0.3),
                  "GOAL", font_size=9, bold=True, color=ACCENT)
    _add_text_box(slide, Inches(0.4), Inches(1.55), Inches(5.5), Inches(0.7),
                  s["goal"], font_size=11, color=TEXT_DARK)

    # Snippet placeholder
    snippet_box = _add_shape_rect(slide, Inches(0.4), Inches(2.4), Inches(5.5), Inches(2.8), CODE_BG)
    tf = snippet_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = s["snippet"]
    p.font.size = Pt(10)
    p.font.name = FONT_CODE
    p.font.color.rgb = TEXT_MED

    # RIGHT COLUMN: Outcome callout
    callout = _add_shape_rect(slide, Inches(6.2), Inches(1.3), Inches(3.4), Inches(2.2),
                               RGBColor(0xEA, 0xF0, 0xF9), ACCENT)
    tf = callout.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "OUTCOME"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = FONT_BODY

    _add_para(tf, s["outcome"], font_size=11, color=TEXT_DARK, space_before=8)

    # Speed-up badge
    speedup_text = s["speedup"]
    badge_color = ACCENT2 if s["speedup_val"] >= 10 else (GRAY if s["speedup_val"] <= 1 else ACCENT)
    speed_badge = _add_shape_rect(slide, Inches(6.4), Inches(3.7), Inches(1.4), Inches(0.5), badge_color)
    tf = speed_badge.text_frame
    p = tf.paragraphs[0]
    p.text = speedup_text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    _add_text_box(slide, Inches(8.0), Inches(3.75), Inches(1.5), Inches(0.4),
                  "speed-up", font_size=10, color=TEXT_MED)

    # Before/After mini display
    _add_text_box(slide, Inches(6.4), Inches(4.35), Inches(3), Inches(0.3),
                  f"{s['before']}  ->  {s['after']}",
                  font_size=12, bold=True, color=TEXT_DARK, font_name=FONT_CODE,
                  alignment=PP_ALIGN.LEFT)

    _add_text_box(slide, Inches(6.4), Inches(4.65), Inches(3), Inches(0.25),
                  f"SAS: {s['sas']}", font_size=9, color=TEXT_LIGHT)

    # Bottom takeaway
    _add_shape_rect(slide, Inches(0.4), Inches(5.5), Inches(9.2), Pt(1), ACCENT)
    tb = _add_text_box(slide, Inches(0.4), Inches(5.6), Inches(9.2), Inches(0.4),
                       s["takeaway"], font_size=11, color=TEXT_MED)
    for p in tb.text_frame.paragraphs:
        p.font.italic = True


def build_takeaways_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
                  "Key Takeaways", font_size=28, bold=True, color=WHITE)

    takeaways = [
        "Avoid functions on indexed/partitioned date columns (WHERE, GROUP BY) — up to 1600x improvement",
        "Tune client fetch sizes (arraysize >= 5000, fetchmany) for bulk reads — 15-16x improvement",
        "Use RESULT_CACHE for repeated read-heavy queries — 22x on subsequent runs",
        "Bind variables reduce parse overhead for looped queries — ~10x for batch lookups",
        "PARALLEL hints limited by SESSIONS_PER_USER=3 in this environment — no gain observed",
        "SAS runtimes closely track Oracle timings — these optimizations benefit both paths",
    ]

    y = Inches(1.4)
    for t in takeaways:
        # Bullet marker
        _add_shape_rect(slide, Inches(0.8), y + Pt(5), Pt(8), Pt(8), ACCENT)
        _add_text_box(slide, Inches(1.2), y, Inches(8), Inches(0.5),
                      t, font_size=13, color=WHITE)
        y += Inches(0.65)


def build_caveats_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_LIGHT)

    _add_text_box(slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.5),
                  "Caveats & Environment", font_size=24, bold=True, color=BG_DARK)

    items = [
        ("Platform", "Oracle via python-oracledb; SAS via SAS/ACCESS passthrough"),
        ("Constraint", "SESSIONS_PER_USER = 3 (limits PARALLEL effectiveness)"),
        ("Methodology", "5 iterations per experiment, averaged; Python timings via time.perf_counter()"),
        ("Tables Tested", "V_TRAN_L (partitioned by EFF_DT), V_PCO_LN_TU_CA, V_LOGON_ATMPT_LOG"),
        ("Disclaimer", "Results are environment-specific; verify optimizations on your target system"),
    ]

    y = Inches(1.3)
    for label, text in items:
        _add_shape_rect(slide, Inches(0.5), y, Inches(1.6), Inches(0.5),
                        RGBColor(0xE3, 0xE8, 0xF0))
        box = _add_text_box(slide, Inches(0.6), y + Pt(4), Inches(1.4), Inches(0.4),
                            label, font_size=10, bold=True, color=BG_DARK,
                            alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, Inches(2.3), y + Pt(2), Inches(7), Inches(0.5),
                      text, font_size=11, color=TEXT_DARK)
        y += Inches(0.7)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_appendix_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)
    _add_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(0.8),
                  "Appendix", font_size=36, bold=True, color=WHITE,
                  alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1), Inches(3), Inches(8), Inches(0.5),
                  "Detailed SQL Optimization Strategies & Results",
                  font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)
    _add_shape_rect(slide, Inches(3.5), Inches(2.85), Inches(3), Pt(3), ACCENT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 widescreen

    # Main slides (1-5)
    build_title_slide(prs)          # 1
    build_summary_slide(prs)        # 2
    build_chart_slide(prs)          # 3
    build_takeaways_slide(prs)      # 4
    build_caveats_slide(prs)        # 5

    # Appendix
    build_appendix_divider(prs)
    for s in STRATEGIES:
        build_strategy_slide(prs, s)

    out = "oracle_benchmark.pptx"
    prs.save(out)
    print(f"Saved {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
