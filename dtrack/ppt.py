"""
PowerPoint generation utilities for dtrack
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from datetime import date
import json
import os
import re
from pathlib import Path
from lxml import etree

# Default color palette
DEFAULT_COLORS = {
    "primary": "#0085CA",
    "dark": "#004990",
    "light_blue": "#87CEFA",
    "callout_bg": "#E8F5FD",
    "white": "#FFFFFF",
    "text": "#333333",
    "light_bg": "#F2F2F2",
}

DEFAULT_FONTS = {
    "title_size": 44,
    "subtitle_size": 24,
    "date_size": 18,
    "section_size": 54,
    "slide_title_size": 32,
    "bullet_size": 18,
    "callout_size": 14,
    "table_header_size": 14,
    "table_body_size": 12,
    "footer_size": 12,
}

DEFAULT_SLIDE = {
    "width": 10,
    "height": 7.5,
    "label": None,
    "label_name": "Restricted - Internal",
}


def _hex_to_rgb(hex_color):
    """Convert hex color string to RGBColor"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def load_config(config_path=None):
    """Load config from JSON file, falling back to defaults"""
    colors = dict(DEFAULT_COLORS)
    fonts = dict(DEFAULT_FONTS)
    slide = dict(DEFAULT_SLIDE)

    if config_path and Path(config_path).exists():
        with open(config_path, encoding='utf-8') as f:
            cfg = json.load(f)
        colors.update(cfg.get("colors", {}))
        fonts.update(cfg.get("fonts", {}))
        slide.update(cfg.get("slide", {}))

    return {
        "colors": {k: _hex_to_rgb(v) for k, v in colors.items()},
        "fonts": fonts,
        "slide": slide,
    }


class PPTBuilder:
    """Builder class for creating PowerPoint presentations"""

    def __init__(self, template_path=None, config_path=None):
        """Initialize with optional template and config"""
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()

        self.cfg = load_config(config_path)
        self.prs.slide_width = Inches(self.cfg["slide"]["width"])
        self.prs.slide_height = Inches(self.cfg["slide"]["height"])
        self.page_num = 1

    @property
    def c(self):
        return self.cfg["colors"]

    @property
    def f(self):
        return self.cfg["fonts"]

    def add_title_slide(self, title, subtitle=None, date_str=None):
        """Add title slide with blue background"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.c["primary"]

        content_box = slide.shapes.add_shape(
            1, Inches(1.5), Inches(2), Inches(7), Inches(3.5)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = self.c["white"]
        content_box.line.fill.background()

        text_frame = content_box.text_frame
        text_frame.vertical_anchor = 1

        p = text_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(self.f["title_size"])
        p.font.bold = True
        p.font.color.rgb = self.c["dark"]

        if subtitle:
            p = text_frame.add_paragraph()
            p.text = subtitle
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(self.f["subtitle_size"])
            p.font.color.rgb = self.c["text"]
            p.space_before = Pt(12)

        if date_str:
            p = text_frame.add_paragraph()
            p.text = date_str
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(self.f["date_size"])
            p.font.color.rgb = self.c["text"]
            p.space_before = Pt(12)

        return slide

    def add_section_slide(self, section_title):
        """Add section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.c["primary"]

        content_box = slide.shapes.add_shape(
            1, Inches(1.5), Inches(1.5), Inches(7), Inches(5)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = self.c["white"]
        content_box.line.fill.background()

        text_frame = content_box.text_frame
        text_frame.vertical_anchor = 1
        text_frame.margin_left = Inches(0.5)

        p = text_frame.paragraphs[0]
        p.text = section_title
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(self.f["section_size"])
        p.font.bold = True
        p.font.color.rgb = self.c["dark"]

        return slide

    def add_content_slide(self, title):
        """Add content slide with standard layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.c["light_bg"]

        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(0.8)
        )
        text_frame = title_box.text_frame
        text_frame.text = title
        p = text_frame.paragraphs[0]
        p.font.size = Pt(self.f["slide_title_size"])
        p.font.bold = True
        p.font.color.rgb = self.c["dark"]

        footer_box = slide.shapes.add_textbox(
            Inches(8.5), Inches(7.2), Inches(1), Inches(0.3)
        )
        text_frame = footer_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = str(self.page_num)
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(self.f["footer_size"])
        p.font.color.rgb = self.c["text"]

        self.page_num += 1
        return slide

    def add_bullets(self, slide, bullets, left=1, top=2, width=8, height=4, font_size=None):
        """Add bullet points to slide"""
        if font_size is None:
            font_size = self.f["bullet_size"]

        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet_text
            p.level = 0
            p.font.size = Pt(font_size)
            p.font.color.rgb = self.c["text"]
            p.space_after = Pt(12)

    def add_callout(self, slide, text, left=1, top=1.8, width=8, height=0.5):
        """Add a colored callout box"""
        shape = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.c["callout_bg"]
        shape.line.color.rgb = self.c["primary"]
        shape.line.width = Pt(1)

        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)

        p = text_frame.paragraphs[0]
        p.font.size = Pt(self.f["callout_size"])
        p.font.color.rgb = self.c["text"]
        p.font.bold = True

    def add_table(self, slide, col_names, rows, left=1.5, top=2, width=7, height=4):
        """Add a table to slide"""
        num_rows = len(rows) + 1
        num_cols = len(col_names)

        table = slide.shapes.add_table(
            num_rows, num_cols, Inches(left), Inches(top), Inches(width), Inches(height)
        ).table

        for col_idx, col_name in enumerate(col_names):
            cell = table.cell(0, col_idx)
            cell.text = col_name
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.c["primary"]

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(self.f["table_header_size"])
            paragraph.font.color.rgb = self.c["white"]

        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_value)

                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(self.f["table_body_size"])
                paragraph.font.color.rgb = self.c["text"]

                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.c["light_bg"]

    def apply_sensitivity_label(self, label_id=None, label_name=None):
        """
        Apply a Microsoft Information Protection (MIP) sensitivity label.
        Falls back to config values if args not provided.

        To find your org's label GUID:
        1. Open a labeled .pptx in PowerPoint
        2. File > Info > Protect Presentation > shows label
        3. Or extract from .pptx: unzip, check docProps/custom.xml
        """
        label_id = label_id or self.cfg["slide"].get("label")
        label_name = label_name or self.cfg["slide"].get("label_name", "Restricted - Internal")

        if not label_id:
            return
        # MIP labels are stored as custom document properties
        # in the docProps/custom.xml part of the .pptx
        CUSTOM_PROPS_URI = "http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"
        VT_URI = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"

        nsmap = {
            None: CUSTOM_PROPS_URI,
            "vt": VT_URI,
        }

        # Build custom.xml with MIP properties
        root = etree.Element(f"{{{CUSTOM_PROPS_URI}}}Properties", nsmap=nsmap)

        mip_props = {
            "MSIP_Label_{}_Enabled": "true",
            "MSIP_Label_{}_SetDate": date.today().isoformat() + "T00:00:00Z",
            "MSIP_Label_{}_Method": "Standard",
            "MSIP_Label_{}_Name": label_name,
            "MSIP_Label_{}_SiteId": "",
            "MSIP_Label_{}_ActionId": "",
            "MSIP_Label_{}_ContentBits": "0",
        }

        for pid, (key_tmpl, value) in enumerate(mip_props.items(), start=2):
            key = key_tmpl.format(label_id)
            prop = etree.SubElement(root, f"{{{CUSTOM_PROPS_URI}}}property")
            prop.set("fmtid", "{D5CDD505-2E9C-101B-9397-08002B2CF9AE}")
            prop.set("pid", str(pid))
            prop.set("name", key)

            vt_lpwstr = etree.SubElement(prop, f"{{{VT_URI}}}lpwstr")
            vt_lpwstr.text = value

        # Add custom properties part to the presentation
        custom_xml = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
        content_type = "application/vnd.openxmlformats-officedocument.custom-properties+xml"

        part = Part(
            PackURI("/docProps/custom.xml"),
            content_type,
            self.prs.part.package,
            blob=custom_xml,
        )

        self.prs.part.package.relate_to(part, RT.CUSTOM_PROPERTIES)

    def save(self, output_path):
        """Save presentation to file. Auto-applies sensitivity label from config."""
        import sys

        # Auto-apply sensitivity label if configured
        self.apply_sensitivity_label()

        output = Path(output_path).resolve()

        try:
            self.prs.save(str(output))
        except PermissionError:
            if sys.platform == 'win32':
                if _pptx_hot_replace(output, self.prs):
                    return str(output)
            # Fallback: save with _latest suffix
            fallback = output.with_stem(f"{output.stem}_latest")
            self.prs.save(str(fallback))
            print(f"  (original is locked, saved to: {fallback.name})")
            return str(fallback)

        return str(output)


def _pptx_hot_replace(file_path, prs):
    """Close file in PowerPoint, save new version, reopen. Requires pywin32."""
    try:
        import win32com.client
    except ImportError:
        print("  tip: pip install pywin32  (enables auto-refresh when PPT is open)")
        return False

    try:
        pp = win32com.client.GetActiveObject("PowerPoint.Application")
    except Exception:
        return False

    # Find matching open presentation
    target = None
    for presentation in pp.Presentations:
        if Path(presentation.FullName).resolve() == file_path:
            target = presentation
            break

    if target is None:
        return False

    # Close → overwrite → reopen
    target.Close()
    prs.save(str(file_path))
    pp.Presentations.Open(str(file_path))
    print("  (PowerPoint refreshed)")
    return True


def parse_markdown_to_ppt(markdown_path, output_path, template_path=None, config_path=None):
    """
    Parse markdown file and create PowerPoint presentation.

    Markdown syntax mapping:
      # Title           → Title slide (first section only)
      ## Subtitle       → Title slide subtitle (first section)
      ### Date          → Title slide date (first section)
      ---               → Slide separator
      ## Heading (alone)→ Section divider slide
      ## Heading + body → Content slide title
      - Item / * Item   → Bullet point
      > Text            → Callout box
      | H | H |         → Table
    """
    builder = PPTBuilder(template_path, config_path)

    with open(markdown_path, 'r', encoding='utf-8') as f:
        content = f.read()

    sections = re.split(r'\n---+\n', content)

    for i, section in enumerate(sections):
        section = section.strip()
        if not section:
            continue

        lines = section.split('\n')

        # Title slide (first section with # heading)
        if i == 0 and lines[0].startswith('# ') and not lines[0].startswith('## '):
            title = lines[0].lstrip('# ').strip()
            subtitle = None
            date_str = None

            for line in lines[1:]:
                if line.startswith('### '):
                    date_str = line.lstrip('# ').strip()
                elif line.startswith('## '):
                    subtitle = line.lstrip('# ').strip()

            builder.add_title_slide(title, subtitle, date_str)
            continue

        # Section divider (only a ## heading, nothing else)
        non_empty = [l for l in lines if l.strip()]
        if len(non_empty) == 1 and non_empty[0].startswith('## '):
            section_title = non_empty[0].lstrip('# ').strip()
            builder.add_section_slide(section_title)
            continue

        # Content slide
        slide_title = None
        bullets = []
        callouts = []
        table_data = []

        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith('## ') and slide_title is None:
                slide_title = stripped.lstrip('# ').strip()
            elif stripped.startswith('> '):
                callouts.append(stripped.lstrip('> ').strip())
            elif stripped.startswith('- ') or stripped.startswith('* '):
                bullets.append(stripped.lstrip('-* ').strip())
            elif '|' in stripped:
                # Skip separator rows
                if re.match(r'^\|[\s\-:|]+\|$', stripped):
                    continue
                cells = [cell.strip() for cell in stripped.split('|')[1:-1]]
                if cells:
                    table_data.append(cells)

        if slide_title:
            slide = builder.add_content_slide(slide_title)

            current_top = 1.8

            # Render callouts first
            for callout_text in callouts:
                builder.add_callout(slide, callout_text, top=current_top)
                current_top += 0.7

            # Render bullets
            if bullets:
                builder.add_bullets(slide, bullets, top=current_top)
                current_top += len(bullets) * 0.5

            # Render tables
            if table_data and len(table_data) > 1:
                headers = table_data[0]
                rows = table_data[1:]
                table_top = max(current_top, 2.0) if not bullets else current_top + 0.5
                builder.add_table(slide, headers, rows, top=table_top)

    builder.save(output_path)
    return output_path
