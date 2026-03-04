#!/usr/bin/env python3
"""
Generate PowerPoint presentation for dtrack project
Barclays-inspired blue theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import date

# Barclays-inspired color palette
BARCLAYS_BLUE = RGBColor(0, 133, 202)  # #0085CA
DARK_BLUE = RGBColor(0, 73, 144)       # #004990
LIGHT_BLUE = RGBColor(135, 206, 250)   # Light sky blue
PALE_BLUE = RGBColor(232, 245, 253)    # Very pale blue
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(242, 242, 242)

def add_header(slide, title_text, subtitle_text=None):
    """Add title and optional subtitle to slide"""
    # Title
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.text = title_text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    if subtitle_text:
        # Add subtitle below title
        left = Inches(1)
        top = Inches(1.3)
        width = Inches(8)
        height = Inches(0.5)

        subtitle = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle.text_frame
        text_frame.text = subtitle_text
        p = text_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY

def add_bullets(slide, left, top, width, height, bullets, font_size=18):
    """Add bullet points to slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

def add_callout(slide, left, top, width, height, text, fill_color=PALE_BLUE):
    """Add a colored callout box"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )

    # Set fill color
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    # Set border
    shape.line.color.rgb = BARCLAYS_BLUE
    shape.line.width = Pt(1)

    # Add text
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)

    p = text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.font.bold = True

def add_side_label(slide, left, top, width, height, label_text):
    """Add blue label box (like Barclays style)"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )

    # Set fill color
    shape.fill.solid()
    shape.fill.fore_color.rgb = BARCLAYS_BLUE

    # No border
    shape.line.fill.background()

    # Add text
    text_frame = shape.text_frame
    text_frame.text = label_text
    text_frame.vertical_anchor = 1  # Middle

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True

def add_table(slide, left, top, width, height, col_names, rows):
    """Add a table to slide"""
    # Calculate table dimensions
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(col_names)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Set column headers
    for col_idx, col_name in enumerate(col_names):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = BARCLAYS_BLUE

        # Format header text
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = WHITE

    # Set row data
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_value)

            # Format data text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = DARK_GRAY

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

def add_footer(slide, page_num):
    """Add footer with page number"""
    left = Inches(8.5)
    top = Inches(7.2)
    width = Inches(1)
    height = Inches(0.3)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = str(page_num)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

def create_title_slide(prs, title, subtitle, date_str):
    """Create title slide with blue background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Blue background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BARCLAYS_BLUE

    # White content box
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(3.5)

    content_box = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = WHITE
    content_box.line.fill.background()

    # Title text
    text_frame = content_box.text_frame
    text_frame.vertical_anchor = 1  # Middle

    p = text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Subtitle
    p = text_frame.add_paragraph()
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(12)

    # Date
    p = text_frame.add_paragraph()
    p.text = date_str
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(12)

def create_section_slide(prs, section_title):
    """Create section divider slide (Barclays style)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Blue background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BARCLAYS_BLUE

    # White content box
    left = Inches(1.5)
    top = Inches(1.5)
    width = Inches(7)
    height = Inches(5)

    content_box = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = WHITE
    content_box.line.fill.background()

    # Section title
    text_frame = content_box.text_frame
    text_frame.vertical_anchor = 1  # Middle

    p = text_frame.paragraphs[0]
    p.text = section_title
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Add left margin
    text_frame.margin_left = Inches(0.5)

def create_content_slide(prs, title, page_num):
    """Create content slide with standard layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Light gray background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY

    # Title bar
    add_header(slide, title)

    # Footer
    add_footer(slide, page_num)

    return slide

def main():
    """Generate the presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    page_num = 1

    # Slide 1: Title
    create_title_slide(
        prs,
        "dtrack: Automated Data Validation & Monitoring",
        "Balance Transfer Tables - Pilot Results",
        "March 4, 2026"
    )

    # Slide 2: Section - Overview
    create_section_slide(prs, "Overview")

    # Slide 3: Executive Summary
    slide = create_content_slide(prs, "Executive Summary", page_num)
    page_num += 1

    # Add labeled items
    labels = ["What", "Why", "Results"]
    contents = [
        "Lightweight CLI tool for tracking data quality over time",
        "Reduce PVT resolution time & increase visibility",
        "9 tables tracked weekly, ready for expansion"
    ]

    for i, (label, content) in enumerate(zip(labels, contents)):
        top_pos = Inches(1.8 + i * 1.2)
        add_side_label(slide, Inches(1), top_pos, Inches(1.5), Inches(0.6), label)
        add_bullets(slide, Inches(2.7), top_pos, Inches(6.3), Inches(0.6), [content], font_size=16)

    # Slide 4: Customer Pain Points
    slide = create_content_slide(prs, "Customer Pain Points", page_num)
    page_num += 1

    add_callout(slide, Inches(1), Inches(1.8), Inches(8), Inches(0.5),
                "Customer questions that drove this:")

    bullets = [
        '"Can I validate data without waiting for complete annotation?"',
        '"What if tables need special processing before comparison?"',
        '"How do I track dozens of tables without scattered emails?"'
    ]
    add_bullets(slide, Inches(1), Inches(2.5), Inches(8), Inches(3), bullets)

    # Slide 5: Section - Current State
    create_section_slide(prs, "Current State")

    # Slide 6: Current Process Pain
    slide = create_content_slide(prs, "Current Process Pain", page_num)
    page_num += 1

    add_callout(slide, Inches(1), Inches(1.8), Inches(8), Inches(0.5),
                "End-to-end execution creates bottlenecks:")

    pain_points = [
        ("Sequential\ndependency", "Business PVT waits for tech PVT + annotation\n(unnecessary delay)"),
        ("Inflexible\nprocessing", "SAS-generated tables or custom logic breaks\nautomation"),
        ("Scope\nlimited", "Validation only at one point in time,\nno ongoing monitoring")
    ]

    for i, (label, description) in enumerate(pain_points):
        top_pos = Inches(2.5 + i * 1.3)
        add_side_label(slide, Inches(1), top_pos, Inches(1.8), Inches(0.8), label)
        add_bullets(slide, Inches(3), top_pos, Inches(6), Inches(0.8), [description], font_size=14)

    # Impact callout
    add_callout(slide, Inches(1), Inches(6.3), Inches(8), Inches(0.6),
                "Impact: Longer resolution time, manual tracking overhead, stakeholder hesitation")

    # Slide 7: Gap Analysis
    slide = create_content_slide(prs, "Gap Analysis", page_num)
    page_num += 1

    col_names = ["Feature", "Current Status"]
    rows = [
        ["Flexible data loading", "❌ Only simple SQL queries"],
        ["Ongoing monitoring", "❌ One-time validation only"],
        ["Central dashboard", "❌ Results scattered in emails"],
        ["Column-level tracking", "⚠️ Partial (row count only)"],
        ["Historical trending", "❌ No time-series view"],
        ["Automated alerts", "❌ Manual review required"]
    ]

    add_table(slide, Inches(1.5), Inches(2), Inches(7), Inches(4), col_names, rows)

    # Slide 8: Section - Solution
    create_section_slide(prs, "Solution")

    # Slide 9: Solution - dtrack
    slide = create_content_slide(prs, "Solution: dtrack", page_num)
    page_num += 1

    add_callout(slide, Inches(1), Inches(1.8), Inches(8), Inches(0.5),
                "Flexible, monitorable data validation:")

    solution_bullets = [
        "Row count & column statistics tracking from any CSV source",
        "Compare paired tables across systems/environments",
        "HTML reports with expandable details & color-coded differences",
        "Weekly updates with historical trending",
        "Handles SAS outputs, custom processing pipelines"
    ]
    add_bullets(slide, Inches(1), Inches(2.5), Inches(8), Inches(4), solution_bullets)

    # Slide 10: Section - Results
    create_section_slide(prs, "Results & Impact")

    # Slide 11: Results & Impact
    slide = create_content_slide(prs, "Results & Impact", page_num)
    page_num += 1

    add_callout(slide, Inches(1), Inches(1.8), Inches(8), Inches(0.5),
                "Pilot: 9 Balance Transfer Tables")

    # Efficiency gains
    gains = [
        ("⏱️ 83% time\nreduction", "2 hours → 20 minutes per validation run"),
        ("📅 2 weeks\nfaster", "Decouple validation from annotation/tech PVT"),
        ("👥 100+ views\nper week", "Central dashboard vs. scattered emails")
    ]

    for i, (label, description) in enumerate(gains):
        top_pos = Inches(2.5 + i * 1.2)
        add_side_label(slide, Inches(1), top_pos, Inches(1.8), Inches(0.7), label)
        add_bullets(slide, Inches(3), top_pos, Inches(6), Inches(0.7), [description], font_size=14)

    # Current status
    status_bullets = [
        "✅ Weekly automated tracking (vs. manual email updates)",
        "✅ Central HTML dashboard accessible to engineers & business",
        "✅ 9 tables monitored with historical trending"
    ]
    add_bullets(slide, Inches(1), Inches(5.8), Inches(8), Inches(1.5), status_bullets, font_size=14)

    # Slide 12: Section - Future
    create_section_slide(prs, "Roadmap")

    # Slide 13: Roadmap
    slide = create_content_slide(prs, "Roadmap", page_num)
    page_num += 1

    # Near-term
    add_callout(slide, Inches(1), Inches(1.8), Inches(8), Inches(0.5),
                "Near-term (Q1-Q2 2026):")

    near_term = [
        "✅ Row count tracking (complete)",
        "🔄 Column-level statistics (in progress)",
        "📋 Automated alerting (planned)"
    ]
    add_bullets(slide, Inches(1), Inches(2.5), Inches(8), Inches(1.5), near_term)

    # Q3 2026 target
    add_callout(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.5),
                "Q3 2026 Target:")

    q3_goals = [
        "📈 Scale to 100+ tables across multiple domains",
        "🔗 Integration with existing BI dashboards",
        "🤖 Anomaly detection & trend analysis"
    ]
    add_bullets(slide, Inches(1), Inches(4.9), Inches(8), Inches(1.8), q3_goals)

    # Slide 14: Appendix
    slide = create_content_slide(prs, "Appendix - Technical Details", page_num)

    # Technical architecture
    add_side_label(slide, Inches(1), Inches(1.8), Inches(2), Inches(0.6), "Architecture")
    arch_bullets = [
        "Database: SQLite with time-series design",
        "Loading: Flexible CSV ingestion from SQL/SAS",
        "Comparison: Paired table validation",
        "Reporting: HTML export with interactive sections"
    ]
    add_bullets(slide, Inches(3.2), Inches(1.8), Inches(5.8), Inches(2.5), arch_bullets, font_size=14)

    # Sample workflow
    add_side_label(slide, Inches(1), Inches(4.5), Inches(2), Inches(0.6), "Workflow")
    workflow_text = """dtrack init project.db
dtrack load-map project.db config.json --type row
dtrack list-pairs project.db -v
dtrack compare-row project.db --pair customer_daily"""

    textbox = slide.shapes.add_textbox(Inches(3.2), Inches(4.5), Inches(5.8), Inches(2))
    text_frame = textbox.text_frame
    text_frame.text = workflow_text
    text_frame.word_wrap = False
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = DARK_GRAY

    # Save presentation
    output_file = "slides/dtrack_presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation saved to: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
